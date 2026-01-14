# generate_ppt.py
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

# 创建演示文稿
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# 设置主题颜色
MAIN_COLOR = RGBColor(46, 134, 171)  # 蓝色
ACCENT_COLOR = RGBColor(230, 57, 70)  # 红色

def add_title_slide(prs, title, subtitle):
    """添加标题页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = MAIN_COLOR
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(54)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # 副标题
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_frame.paragraphs[0].font.size = Pt(28)
    subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, image_path, content_points=None):
    """添加内容页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # 标题栏
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(40)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = MAIN_COLOR
    
    # 图像
    if image_path and os.path.exists(f'ppt_figures/{image_path}'):
        slide.shapes.add_picture(f'ppt_figures/{image_path}', 
                                Inches(0.5), Inches(1.1), width=Inches(9))
    
    # 内容点
    if content_points:
        y_pos = 6
        for point in content_points:
            text_box = slide.shapes.add_textbox(Inches(0.7), Inches(y_pos), Inches(8.6), Inches(0.4))
            text_frame = text_box.text_frame
            text_frame.text = f"• {point}"
            text_frame.paragraphs[0].font.size = Pt(14)
            y_pos += 0.35

# 创建幻灯片
add_title_slide(prs, "GPU并行DSP信号处理系统", "高性能实时信号处理解决方案")

add_content_slide(prs, "系统架构",  "01_architecture.png", [
    "✓ 模块化设计：数据预处理 → FFT → 滤波 → 脉冲压缩",
    "✓ GPU优化层：显存管理、CUDA Streams、Kernel融合、流式处理",
    "✓ 高效执行：批处理、异步传输、动态调度"
])

add_content_slide(prs, "性能指标", "02_performance_comparison.png", [
    "✓ FFT加速: 83.8x - 115.8x",
    "✓ 脉冲压缩: 38.2x - 104x",
    "✓ 流式处理: 55.5 Msps",
    "✓ 并行度: 1.40x"
])

add_content_slide(prs, "处理流程", "03_processing_pipeline.png", [
    "✓ 9个处理阶段，集成GPU优化",
    "✓ 支持实时流式处理",
    "✓ <10ms端到端延迟"
])

add_content_slide(prs, "GPU部署架构", "04_gpu_deployment.png", [
    "✓ RTX/A100 GPU部署",
    "✓ 多流并行执行",
    "✓ 显存与PCIe优化"
])

add_content_slide(prs, "数据流与内存管理", "05_data_flow_memory.png", [
    "✓ 显存分层利用",
    "✓ Pinned内存优化",
    "✓ 异步H2D/D2H传输"
])

add_content_slide(prs, "优化策略分析", "06_optimization_analysis.png", [
    "✓ 批处理: 2.75x加速",
    "✓ Kernel融合: 1.16x加速",
    "✓ CUDA Streams: 1.40x加速",
    "✓ 显存预分配: 1.12x加速"
])

# 最后一页：总结
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(42, 157, 143)

title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
title_frame = title_box.text_frame
title_frame.text = "项目成果"
title_frame.paragraphs[0].font.size = Pt(54)
title_frame.paragraphs[0].font.bold = True
title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

content_box = slide.shapes.add_textbox(Inches(1.5), Inches(4), Inches(7), Inches(2.5))
content_frame = content_box.text_frame
content_frame.word_wrap = True

achievements = [
    "✓ 100x+性能加速",
    "✓ 55.5 Msps流式处理能力",
    "✓ <10ms实时响应延迟",
    "✓ 生产级系统部署"
]

for i, achievement in enumerate(achievements):
    if i == 0:
        content_frame.text = achievement
        p = content_frame.paragraphs[0]
    else:
        p = content_frame.add_paragraph()
        p.text = achievement
        p.level = 0
    
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.space_before = Pt(12)

# 保存
prs.save('GPU_DSP系统_汇报.pptx')
print("✓ PPT已生成: GPU_DSP系统_汇报.pptx")
